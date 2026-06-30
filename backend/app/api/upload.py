from fastapi import APIRouter, UploadFile, HTTPException
from sqlalchemy.orm import Session
from app.db.session import SessionLocal
from app.core.config import settings
from langchain_openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.schema import Document
import os
import fitz  # PyMuPDF
from docx import Document as DocxDocument
from pptx import Presentation

router = APIRouter()

def extract_text_from_pdf(path: str) -> str:
    with fitz.open(path) as doc:
        return "".join(page.get_text() for page in doc)

def extract_text_from_docx(path: str) -> str:
    doc = DocxDocument(path)
    return "\n".join(p.text for p in doc.paragraphs)

def extract_text_from_pptx(path: str) -> str:
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

@router.post("/upload/")
async def upload(file: UploadFile):
    db: Session = SessionLocal()

    user_id = 1
    topic_id = 1

    ext = os.path.splitext(file.filename)[1].lower()
    allowed_extensions = [".pdf", ".docx", ".pptx"]
    if ext not in allowed_extensions:
        raise HTTPException(status_code=400, detail="지원하지 않는 파일 형식입니다. (pdf, docx, pptx만 가능)")

    # 1. 파일 저장
    contents = await file.read()
    os.makedirs("./uploads", exist_ok=True)
    file_path = f"./uploads/{file.filename}"
    with open(file_path, "wb") as f:
        f.write(contents)

    # 2. 텍스트 추출
    try:
        if ext == ".pdf":
            raw_text = extract_text_from_pdf(file_path)
        elif ext == ".docx":
            raw_text = extract_text_from_docx(file_path)
        elif ext == ".pptx":
            raw_text = extract_text_from_pptx(file_path)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"텍스트 추출 실패: {str(e)}")

    if not raw_text.strip():
        raise HTTPException(status_code=400, detail="문서에 추출할 텍스트가 없습니다.")

    # 3. 텍스트 분할
    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    chunks = splitter.split_text(raw_text)

    # text_splitter = RecursiveCharacterTextSplitter.from_tiktoken_encoder(
    #     model_name="gpt-4",
    #     chunk_size=100,
    #     chunk_overlap=0,
    # )

    documents = [
        Document(page_content=chunk, metadata={"document_id": topic_id, "chunk_index": i})
        for i, chunk in enumerate(chunks)
    ]

    # 4. 벡터 저장
    embeddings = OpenAIEmbeddings(model="text-embedding-3-small", openai_api_key=settings.openai_api_key)
    faiss_path = "./faiss_store"
    if not os.path.exists(faiss_path):
        vs = FAISS.from_documents(documents, embeddings)
    else:
        vs = FAISS.load_local(faiss_path, embeddings)
        vs.add_documents(documents)
    vs.save_local(faiss_path)

    return {"status": f"{file.filename} 업로드 및 벡터 저장 완료"}