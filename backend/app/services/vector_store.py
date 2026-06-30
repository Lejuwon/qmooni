import os
import fitz
from typing import List
from docx import Document as DocxDocument
from pptx import Presentation
# from langchain_openai import OpenAIEmbeddings
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain_community.vectorstores import FAISS
from langchain.schema import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
from app.core.config import settings
import numpy as np
from app.services.quiz_generator_mcq import summarize_document

# 1. 텍스트 추출
def extract_text(file_path: str, filename: str) -> List[Document]:
    ext = os.path.splitext(filename)[1].lower()
    documents = []

    if ext == ".pdf":
        with fitz.open(file_path) as doc:
            for i, page in enumerate(doc):
                text = page.get_text()
                if text.strip():
                    documents.append(Document(page_content=text, metadata={"page": i + 1}))
        return documents

    elif ext == ".docx":
        doc = DocxDocument(file_path)
        for i, para in enumerate(doc.paragraphs):
            text = para.text.strip()
            if text:
                documents.append(Document(page_content=text, metadata={"page": i + 1}))
        return documents

    elif ext == ".pptx":
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            text = "\n".join(
                shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()
            )
            if text.strip():
                documents.append(Document(page_content=text, metadata={"page": i + 1}))
        return documents

    else:
        raise ValueError("지원하지 않는 형식입니다.")

# 2. 문서 벡터화 후 저장
def process_and_embed(file_path: str, filename: str, user_id: int, topic_id: int, document_id: int):
    raw_docs = extract_text(file_path, filename)
    if not raw_docs:
        raise ValueError("문서에 추출할 텍스트가 없습니다.")

    splitter = RecursiveCharacterTextSplitter(chunk_size=500, chunk_overlap=50)
    all_chunks = []

    for i, doc in enumerate(raw_docs):
        # ✅ 1. 원문 청크 저장
        chunks = splitter.split_text(doc.page_content)
        for j, chunk in enumerate(chunks):
            all_chunks.append(Document(
                page_content=chunk,
                metadata={
                    "user_id": user_id,
                    "topic_id": topic_id,
                    "document_id": document_id,
                    "chunk_index": len(all_chunks),
                    "page": doc.metadata.get("page"),
                    "file_name": filename,
                    "type": "raw"  # 원문 청크
                }
            ))

        # ✅ 2. 페이지 요약 생성
        try:
            summary = summarize_document([doc.page_content])  # 사용자가 정의한 GPT 요약 함수
            all_chunks.append(Document(
                page_content=summary,
                metadata={
                    "user_id": user_id,
                    "topic_id": topic_id,
                    "document_id": document_id,
                    "chunk_index": len(all_chunks),
                    "page": doc.metadata.get("page"),
                    "file_name": filename,
                    "type": "summary"  # 요약 청크
                }
            ))
        except Exception as e:
            print(f"⚠️ 페이지 {doc.metadata.get('page')} 요약 실패: {e}")

    # ✅ FAISS에 저장
    faiss_path = f"./faiss_store/user_{user_id}/topic_{topic_id}"
    os.makedirs(faiss_path, exist_ok=True)

    embeddings = OpenAIEmbeddings(model="text-embedding-3-small", openai_api_key=settings.openai_api_key)

    if not os.path.exists(os.path.join(faiss_path, "index.faiss")):
        vs = FAISS.from_documents(all_chunks, embeddings)
    else:
        vs = FAISS.load_local(faiss_path, embeddings, allow_dangerous_deserialization=True)
        vs.add_documents(all_chunks)

    vs.save_local(faiss_path)


def remove_doc_vs(user_id: int, topic_id: int, document_id: int):
    faiss_path = f"./faiss_store/user_{user_id}/topic_{topic_id}"
    index_path = os.path.join(faiss_path, "index.faiss")
    if not os.path.exists(index_path):
        return  # 벡터가 없는 경우는 조용히 무시

    embeddings = OpenAIEmbeddings(model="text-embedding-3-small", openai_api_key=settings.openai_api_key)
    vs = FAISS.load_local(faiss_path, embeddings, allow_dangerous_deserialization=True)

    docstore = vs.docstore._dict
    remove_keys = [k for k, v in docstore.items() if v.metadata.get("document_id") == document_id]

    for key in remove_keys:
        del docstore[key]
        vs.index.remove_ids(np.array([key], dtype='object'))  # ← 핵심 수정!

    vs.save_local(faiss_path)